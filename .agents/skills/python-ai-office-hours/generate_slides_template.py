"""Generate a single-slide PPTX for Python+AI Office Hours weekly updates (light mode).

Template for the python-ai-office-hours Copilot CLI skill.
Replace the placeholder bullets with the current week's content.
Update OUTPUT_DATE to the current Monday's date.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import date

# ── Configuration ──
OUTPUT_DATE = date.today().strftime("%B %-d, %Y")  # e.g. "April 7, 2026"
OUTPUT_FILE = f"office-hours-{date.today().isoformat()}.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# -- Light mode palette --
BG = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
LINK = RGBColor(0x6C, 0x75, 0x7D)
TITLE_COLOR = RGBColor(0x0B, 0x5E, 0xD7)
MS_COLOR = RGBColor(0x6F, 0x2D, 0xA8)    # purple — Microsoft/GitHub
IND_COLOR = RGBColor(0x1A, 0x7F, 0x37)   # green — Industry
MY_COLOR = RGBColor(0x0D, 0x6E, 0xFD)    # blue — My work

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = BG

# ── Helpers ──
def add_line(tf, text, size=9, color=BLACK, bold=False, space_after=Pt(1)):
    p = tf.add_paragraph()
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def section_header(tf, text, color=MS_COLOR):
    add_line(tf, text, size=16, color=color, bold=True, space_after=Pt(4))

def bullet(tf, text, link=None):
    p = tf.add_paragraph()
    p.space_after = Pt(3)
    run = p.add_run()
    run.text = f"• {text}"
    run.font.size = Pt(13)
    run.font.color.rgb = BLACK
    if link:
        r2 = p.add_run()
        r2.text = f"  {link}"
        r2.font.size = Pt(10)
        r2.font.color.rgb = LINK
        r2.hyperlink.address = link if link.startswith("http") else f"https://{link}"

# ── Title bar ──
title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.7))
title_box.text_frame.word_wrap = True
ttf = title_box.text_frame
ttf.paragraphs[0].text = ""
add_line(ttf, f"🐍 Python + AI Office Hours — Weekly News Roundup — {OUTPUT_DATE}",
         size=22, color=TITLE_COLOR, bold=True, space_after=Pt(0))

# ── Left column: Microsoft / GitHub ──
left = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(4.1), Inches(6.4))
left.text_frame.word_wrap = True
ltf = left.text_frame
ltf.paragraphs[0].text = ""

section_header(ltf, "🏢 Microsoft / GitHub", MS_COLOR)
# === REPLACE these bullets with this week's MS/GitHub news ===
bullet(ltf, "Example: New Azure OpenAI feature launched", "learn.microsoft.com")
bullet(ltf, "Example: GitHub Copilot update", "github.blog")

# ── Middle column: Industry ──
mid = slide.shapes.add_textbox(Inches(4.7), Inches(0.9), Inches(4.1), Inches(6.4))
mid.text_frame.word_wrap = True
mtf = mid.text_frame
mtf.paragraphs[0].text = ""

section_header(mtf, "🌐 Industry", IND_COLOR)
# === REPLACE these bullets with this week's industry news ===
bullet(mtf, "Example: New open-source AI tool released", "github.com/example")
bullet(mtf, "Example: Python 3.x update", "python.org")

# ── Right column: My Work ──
right = slide.shapes.add_textbox(Inches(9.0), Inches(0.9), Inches(4.0), Inches(6.4))
right.text_frame.word_wrap = True
rtf = right.text_frame
rtf.paragraphs[0].text = ""

section_header(rtf, "👩‍💻 What I've Been Up To", MY_COLOR)
# === REPLACE these bullets with Pamela's activity this week ===
bullet(rtf, "Example: Blog post about MCP", "blog.pamelafox.org")
bullet(rtf, "Example: New PR for Responses API migration")

# ── Save ──
prs.save(OUTPUT_FILE)
print(f"✅ Saved to {OUTPUT_FILE}")
